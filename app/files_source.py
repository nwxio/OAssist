from __future__ import annotations

import csv
import difflib
import hashlib
import json
import mimetypes
import os
import re
import secrets
import shutil
import sqlite3
import uuid
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Any

from app.config import Settings
from app.document_ids import SOURCE_FILES, make_document_id
from app.rag_indexing import delete_document_from_rag, upsert_document_to_rag

FILES_ACCESS_READ_ONLY = "read-only"
FILES_ACCESS_READ_WRITE = "read-write"
VALID_FILES_ACCESS_MODES = {FILES_ACCESS_READ_ONLY, FILES_ACCESS_READ_WRITE}

DEFAULT_FILES_FEATURE_SETTINGS: dict[str, Any] = {
    "access_mode": FILES_ACCESS_READ_ONLY,
    "upload_max_size_mb": 25,
    "upload_ttl_hours": 48,
    "upload_allow_indexable": True,
    "upload_allowed_extensions": [],
    "artifact_ttl_days": 30,
    "artifact_max_size_mb": 20,
}

TEXT_INDEX_EXTENSIONS = {
    ".txt",
    ".md",
    ".rst",
    ".csv",
    ".tsv",
    ".json",
    ".yaml",
    ".yml",
    ".xml",
    ".ini",
    ".conf",
    ".log",
    ".html",
    ".css",
    ".scss",
    ".less",
    ".js",
    ".mjs",
    ".cjs",
    ".ts",
    ".jsx",
    ".tsx",
    ".py",
    ".java",
    ".kt",
    ".go",
    ".rs",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cs",
    ".php",
    ".rb",
    ".swift",
    ".sql",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
}

OFFICE_INDEX_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
INDEXABLE_EXTENSIONS = TEXT_INDEX_EXTENSIONS.union(OFFICE_INDEX_EXTENSIONS)

TEXT_WRITE_EXTENSIONS = {
    ".txt",
    ".md",
    ".rst",
    ".csv",
    ".tsv",
    ".json",
    ".yaml",
    ".yml",
    ".xml",
    ".ini",
    ".conf",
    ".log",
    ".html",
    ".css",
    ".scss",
    ".less",
    ".js",
    ".mjs",
    ".cjs",
    ".ts",
    ".jsx",
    ".tsx",
    ".py",
    ".java",
    ".kt",
    ".go",
    ".rs",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cs",
    ".php",
    ".rb",
    ".swift",
    ".sql",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
    ".dockerfile",
    ".makefile",
}


def _connect(db_path: str) -> sqlite3.Connection:
    conn = sqlite3.connect(db_path, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    return conn


def _utc_now() -> str:
    return datetime.now(UTC).isoformat(timespec="seconds")


def _utc_now_precise() -> str:
    return datetime.now(UTC).isoformat(timespec="microseconds")


def _parse_iso8601(value: str | None) -> datetime | None:
    text = str(value or "").strip()
    if not text:
        return None
    try:
        parsed = datetime.fromisoformat(text)
    except ValueError:
        return None
    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=UTC)
    return parsed.astimezone(UTC)


def _ensure_directory(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def _sanitize_name(filename: str, fallback: str = "file") -> str:
    text = str(filename or "").strip().replace("\\", "/")
    text = text.split("/")[-1]
    text = re.sub(r"[^A-Za-z0-9._-]+", "_", text)
    text = text.strip("._")
    if not text:
        text = fallback
    if len(text) > 120:
        stem, suffix = os.path.splitext(text)
        text = f"{stem[:96]}{suffix[:24]}"
    return text


def _normalize_upload_filename(filename: str, fallback: str = "upload") -> str:
    text = str(filename or "").strip().replace("\\", "/")
    text = text.split("/")[-1]
    text = re.sub(r"[\x00-\x1f\x7f]+", " ", text).strip()
    if not text or text in {".", ".."}:
        text = fallback
    if len(text) > 180:
        stem, suffix = os.path.splitext(text)
        text = f"{stem[:148]}{suffix[:32]}"
    return text


def _next_available_filename(directory: Path, preferred_name: str) -> str:
    base_name = _normalize_upload_filename(preferred_name, fallback="upload")
    stem = Path(base_name).stem or "upload"
    suffix = Path(base_name).suffix

    candidate = base_name
    index = 1
    while (directory / candidate).exists():
        candidate = f"{stem} ({index}){suffix}" if suffix else f"{stem} ({index})"
        index += 1
        if index > 9999:
            candidate = f"{stem}-{uuid.uuid4().hex[:8]}{suffix}" if suffix else f"{stem}-{uuid.uuid4().hex[:8]}"
            if not (directory / candidate).exists():
                break
    return candidate


def _normalize_rel_path(path: str) -> str:
    text = str(path or "").strip().replace("\\", "/")
    text = re.sub(r"/{2,}", "/", text)
    text = text.strip("/")
    if text in {"", "."}:
        return ""
    parts = []
    for part in text.split("/"):
        part = part.strip()
        if not part or part in {".", ".."}:
            continue
        parts.append(part)
    return "/".join(parts)


def _safe_path_join(root: Path, rel_path: str) -> Path:
    candidate = (root / _normalize_rel_path(rel_path)).resolve()
    root_real = root.resolve()
    if candidate != root_real and root_real not in candidate.parents:
        raise ValueError("Path escapes configured root")
    return candidate


def _suffixes(filename: str) -> set[str]:
    name = str(filename or "").strip().lower()
    suffix = Path(name).suffix.lower()
    result = {suffix} if suffix else set()
    if name == "dockerfile":
        result.add(".dockerfile")
    if name == "makefile":
        result.add(".makefile")
    return result


def _is_ext_allowed(filename: str, allowed_extensions: set[str]) -> bool:
    if not allowed_extensions:
        return True
    suffixes = _suffixes(filename)
    return bool(suffixes.intersection(allowed_extensions))


def ensure_files_feature_schema(settings: Settings) -> None:
    data_root = Path("/data")
    _ensure_directory(data_root)
    _ensure_directory(data_root / "files_source")
    _ensure_directory(data_root / "chat_uploads")
    _ensure_directory(data_root / "chat_artifacts")

    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS files_roots (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              name TEXT NOT NULL,
              root_path TEXT NOT NULL UNIQUE,
              enabled INTEGER NOT NULL DEFAULT 1,
              created_at TEXT NOT NULL,
              updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS files_settings (
              key TEXT PRIMARY KEY,
              value TEXT NOT NULL,
              updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS files_documents (
              document_id TEXT PRIMARY KEY,
              root_id INTEGER NOT NULL,
              rel_path TEXT NOT NULL,
              mtime TEXT NOT NULL,
              size INTEGER NOT NULL,
              hash TEXT NOT NULL,
              indexed_at TEXT NOT NULL,
              FOREIGN KEY(root_id) REFERENCES files_roots(id) ON DELETE CASCADE
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_uploads (
              id TEXT PRIMARY KEY,
              user_id INTEGER NOT NULL,
              chat_id TEXT,
              mode TEXT NOT NULL,
              mime TEXT,
              size INTEGER NOT NULL,
              filename TEXT NOT NULL,
              storage_path TEXT NOT NULL,
              indexed_document_id TEXT,
              extracted_text TEXT,
              expires_at TEXT NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_artifacts (
              id TEXT PRIMARY KEY,
              user_id INTEGER NOT NULL,
              chat_id TEXT,
              format TEXT NOT NULL,
              filename TEXT NOT NULL,
              storage_path TEXT NOT NULL,
              size INTEGER NOT NULL,
              sha256 TEXT NOT NULL,
              expires_at TEXT NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_artifact_tokens (
              token TEXT PRIMARY KEY,
              artifact_id TEXT NOT NULL,
              expires_at TEXT NOT NULL,
              used_at TEXT,
              created_at TEXT NOT NULL,
              FOREIGN KEY(artifact_id) REFERENCES chat_artifacts(id) ON DELETE CASCADE
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS files_write_audit (
              id TEXT PRIMARY KEY,
              user_id INTEGER NOT NULL,
              op TEXT NOT NULL,
              path TEXT,
              document_id TEXT,
              before_hash TEXT,
              after_hash TEXT,
              status TEXT NOT NULL,
              message TEXT,
              created_at TEXT NOT NULL,
              FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            )
            """
        )

        conn.execute("CREATE INDEX IF NOT EXISTS idx_files_documents_root_id ON files_documents(root_id)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_uploads_user_created ON chat_uploads(user_id, created_at DESC)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_uploads_expires ON chat_uploads(expires_at)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_artifacts_user_created ON chat_artifacts(user_id, created_at DESC)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_artifacts_expires ON chat_artifacts(expires_at)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_artifact_tokens_artifact ON chat_artifact_tokens(artifact_id)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_files_write_audit_user_created ON files_write_audit(user_id, created_at DESC)")

        now = _utc_now()
        conn.executemany(
            """
            INSERT INTO files_settings (key, value, updated_at)
            VALUES (?, ?, ?)
            ON CONFLICT(key) DO NOTHING
            """,
            [
                (key, json.dumps(value, ensure_ascii=False), now)
                for key, value in DEFAULT_FILES_FEATURE_SETTINGS.items()
            ],
        )
        conn.commit()


def _load_files_settings_raw(settings: Settings) -> dict[str, Any]:
    state = dict(DEFAULT_FILES_FEATURE_SETTINGS)
    with _connect(settings.auth_db_path) as conn:
        rows = conn.execute("SELECT key, value FROM files_settings").fetchall()
    for row in rows:
        key = str(row["key"] or "").strip()
        if key not in state:
            continue
        try:
            parsed = json.loads(str(row["value"] or "null"))
        except json.JSONDecodeError:
            continue
        state[key] = parsed
    return state


def get_files_feature_settings(settings: Settings) -> dict[str, Any]:
    state = _load_files_settings_raw(settings)
    access_mode = str(state.get("access_mode") or FILES_ACCESS_READ_ONLY).strip().lower()
    if access_mode not in VALID_FILES_ACCESS_MODES:
        access_mode = FILES_ACCESS_READ_ONLY

    upload_max_size_mb = int(state.get("upload_max_size_mb") or DEFAULT_FILES_FEATURE_SETTINGS["upload_max_size_mb"])
    upload_ttl_hours = int(state.get("upload_ttl_hours") or DEFAULT_FILES_FEATURE_SETTINGS["upload_ttl_hours"])
    upload_allow_indexable = bool(state.get("upload_allow_indexable", True))

    allowed_ext_raw = state.get("upload_allowed_extensions")
    if isinstance(allowed_ext_raw, list):
        allowed = {str(item).strip().lower() for item in allowed_ext_raw if str(item).strip()}
    else:
        allowed = set()
    allowed = {item if item.startswith(".") else f".{item}" for item in allowed}

    artifact_ttl_days = int(state.get("artifact_ttl_days") or DEFAULT_FILES_FEATURE_SETTINGS["artifact_ttl_days"])
    artifact_max_size_mb = int(state.get("artifact_max_size_mb") or DEFAULT_FILES_FEATURE_SETTINGS["artifact_max_size_mb"])

    return {
        "access_mode": access_mode,
        "upload_max_size_mb": max(1, min(upload_max_size_mb, 500)),
        "upload_ttl_hours": max(1, min(upload_ttl_hours, 24 * 30)),
        "upload_allow_indexable": upload_allow_indexable,
        "upload_allowed_extensions": sorted(allowed),
        "artifact_ttl_days": max(1, min(artifact_ttl_days, 365)),
        "artifact_max_size_mb": max(1, min(artifact_max_size_mb, 200)),
    }


def set_files_feature_settings(settings: Settings, patch: dict[str, Any]) -> dict[str, Any]:
    allowed_keys = set(DEFAULT_FILES_FEATURE_SETTINGS.keys())
    updates: dict[str, Any] = {}
    for key, value in (patch or {}).items():
        if key not in allowed_keys:
            continue
        updates[key] = value

    if "access_mode" in updates:
        mode = str(updates["access_mode"] or "").strip().lower()
        if mode not in VALID_FILES_ACCESS_MODES:
            raise ValueError("access_mode must be 'read-only' or 'read-write'")
        updates["access_mode"] = mode

    if "upload_allowed_extensions" in updates:
        raw = updates.get("upload_allowed_extensions")
        if not isinstance(raw, list):
            raise ValueError("upload_allowed_extensions must be an array")
        cleaned = []
        for item in raw:
            text = str(item or "").strip().lower()
            if not text:
                continue
            cleaned.append(text if text.startswith(".") else f".{text}")
        updates["upload_allowed_extensions"] = sorted(set(cleaned))

    now = _utc_now()
    with _connect(settings.auth_db_path) as conn:
        for key, value in updates.items():
            conn.execute(
                """
                INSERT INTO files_settings (key, value, updated_at)
                VALUES (?, ?, ?)
                ON CONFLICT(key) DO UPDATE SET
                  value=excluded.value,
                  updated_at=excluded.updated_at
                """,
                (key, json.dumps(value, ensure_ascii=False), now),
            )
        conn.commit()
    return get_files_feature_settings(settings)


def files_access_mode(settings: Settings) -> str:
    return str(get_files_feature_settings(settings).get("access_mode") or FILES_ACCESS_READ_ONLY)


def list_files_roots(settings: Settings, *, include_disabled: bool = True) -> list[dict[str, Any]]:
    query = "SELECT id, name, root_path, enabled, created_at, updated_at FROM files_roots"
    params: tuple[Any, ...] = ()
    if not include_disabled:
        query += " WHERE enabled=1"
    query += " ORDER BY name COLLATE NOCASE ASC, id ASC"
    with _connect(settings.auth_db_path) as conn:
        rows = conn.execute(query, params).fetchall()

    roots: list[dict[str, Any]] = []
    for row in rows:
        root_path = str(row["root_path"] or "")
        root_obj = Path(root_path)
        roots.append(
            {
                "id": int(row["id"]),
                "name": str(row["name"] or ""),
                "root_path": root_path,
                "enabled": bool(int(row["enabled"])),
                "exists": root_obj.exists(),
                "is_dir": root_obj.is_dir(),
                "created_at": str(row["created_at"]),
                "updated_at": str(row["updated_at"]),
            }
        )
    return roots


def _canonical_root_path(raw_path: str) -> str:
    text = str(raw_path or "").strip()
    if not text:
        raise ValueError("root_path is required")
    path = Path(text).expanduser()
    if not path.is_absolute():
        raise ValueError("root_path must be absolute")
    resolved = path.resolve(strict=False)
    if not resolved.exists():
        raise ValueError("root_path directory does not exist inside container")
    if not resolved.is_dir():
        raise ValueError("root_path must point to a directory")
    return str(resolved)


def create_files_root(settings: Settings, *, name: str, root_path: str, enabled: bool = True) -> dict[str, Any]:
    clean_name = str(name or "").strip()
    if not clean_name:
        raise ValueError("name is required")
    clean_path = _canonical_root_path(root_path)
    now = _utc_now()
    with _connect(settings.auth_db_path) as conn:
        try:
            cur = conn.execute(
                "INSERT INTO files_roots (name, root_path, enabled, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                (clean_name, clean_path, 1 if enabled else 0, now, now),
            )
        except sqlite3.IntegrityError as exc:
            raise ValueError("root_path already exists") from exc
        conn.commit()
        last_id = cur.lastrowid
        if last_id is None:
            raise ValueError("failed to create root")
        root_id = int(last_id)
    row = get_files_root(settings, root_id)
    if row is None:
        raise ValueError("failed to create root")
    return row


def get_files_root(settings: Settings, root_id: int) -> dict[str, Any] | None:
    with _connect(settings.auth_db_path) as conn:
        row = conn.execute(
            "SELECT id, name, root_path, enabled, created_at, updated_at FROM files_roots WHERE id=?",
            (int(root_id),),
        ).fetchone()
    if row is None:
        return None
    root_path = str(row["root_path"] or "")
    root_obj = Path(root_path)
    return {
        "id": int(row["id"]),
        "name": str(row["name"] or ""),
        "root_path": root_path,
        "enabled": bool(int(row["enabled"])),
        "exists": root_obj.exists(),
        "is_dir": root_obj.is_dir(),
        "created_at": str(row["created_at"]),
        "updated_at": str(row["updated_at"]),
    }


def update_files_root(
    settings: Settings,
    root_id: int,
    *,
    name: str | None = None,
    root_path: str | None = None,
    enabled: bool | None = None,
) -> dict[str, Any]:
    current = get_files_root(settings, root_id)
    if current is None:
        raise ValueError("root not found")

    new_name = str(name).strip() if name is not None else str(current["name"])
    if not new_name:
        raise ValueError("name cannot be empty")

    new_path = _canonical_root_path(root_path) if root_path is not None else str(current["root_path"])
    new_enabled = bool(enabled) if enabled is not None else bool(current["enabled"])
    now = _utc_now()

    with _connect(settings.auth_db_path) as conn:
        try:
            conn.execute(
                "UPDATE files_roots SET name=?, root_path=?, enabled=?, updated_at=? WHERE id=?",
                (new_name, new_path, 1 if new_enabled else 0, now, int(root_id)),
            )
        except sqlite3.IntegrityError as exc:
            raise ValueError("root_path already exists") from exc
        conn.commit()

    updated = get_files_root(settings, root_id)
    if updated is None:
        raise ValueError("root not found after update")
    return updated


def delete_files_root(settings: Settings, root_id: int) -> None:
    with _connect(settings.auth_db_path) as conn:
        conn.execute("DELETE FROM files_roots WHERE id=?", (int(root_id),))
        conn.execute("DELETE FROM files_documents WHERE root_id=?", (int(root_id),))
        conn.commit()


def files_document_id(root_id: int, rel_path: str) -> str:
    normalized = _normalize_rel_path(rel_path)
    if not normalized:
        return make_document_id(SOURCE_FILES, f"{int(root_id)}:/")
    return make_document_id(SOURCE_FILES, f"{int(root_id)}:{normalized}")


def files_folder_document_id(root_id: int, rel_path: str) -> str:
    normalized = _normalize_rel_path(rel_path)
    if not normalized:
        return make_document_id(SOURCE_FILES, f"{int(root_id)}:/")
    return make_document_id(SOURCE_FILES, f"{int(root_id)}:{normalized}/")


def parse_files_document_id(document_id: str) -> tuple[int, str, bool] | None:
    text = str(document_id or "").strip()
    if not text.lower().startswith("files:"):
        return None
    native = text[6:]
    if ":" not in native:
        return None
    root_part, rel = native.split(":", 1)
    try:
        root_id = int(root_part)
    except ValueError:
        return None
    is_folder = rel.endswith("/")
    if rel == "/":
        return root_id, "", True
    rel_path = _normalize_rel_path(rel.rstrip("/"))
    return root_id, rel_path, is_folder


def _walk_root_files(root_path: Path) -> list[str]:
    rel_paths: list[str] = []
    for dirpath, dirnames, filenames in os.walk(root_path, topdown=True):
        dirnames[:] = [name for name in dirnames if name not in {".git", ".trash", "__pycache__"}]
        for filename in filenames:
            if filename.startswith("."):
                continue
            full = Path(dirpath) / filename
            try:
                rel = str(full.relative_to(root_path)).replace("\\", "/")
            except ValueError:
                continue
            rel_paths.append(rel)
    rel_paths.sort()
    return rel_paths


def _file_hash_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _read_text_fallback(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _clean_html(text: str) -> str:
    value = re.sub(r"<script[\s\S]*?</script>", " ", text, flags=re.IGNORECASE)
    value = re.sub(r"<style[\s\S]*?</style>", " ", value, flags=re.IGNORECASE)
    value = re.sub(r"<(nav|footer|header)[\s\S]*?</\1>", " ", value, flags=re.IGNORECASE)
    value = re.sub(r"<[^>]+>", " ", value)
    value = re.sub(r"\s+", " ", value).strip()
    return value


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except Exception:
        return ""
    try:
        doc = Document(str(path))
    except Exception:
        return ""
    parts = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text and paragraph.text.strip()]
    return "\n".join(parts).strip()


def _extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return ""
    try:
        reader = PdfReader(str(path))
    except Exception:
        return ""
    lines: list[str] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if text:
            lines.append(text)
    return "\n\n".join(lines).strip()


def _extract_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception:
        return ""
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return ""
    chunks: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"# Sheet: {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(chunks).strip()


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        presentation = Presentation(str(path))
    except Exception:
        return ""
    slides: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                blocks.append(text.strip())
        if blocks:
            slides.append(f"# Slide {idx}\n" + "\n".join(blocks))
    return "\n\n".join(slides).strip()


def extract_file_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if path.name.lower() in {"dockerfile", "makefile"}:
        suffix = f".{path.name.lower()}"

    try:
        if suffix in {".csv", ".tsv"}:
            delimiter = "," if suffix == ".csv" else "\t"
            lines: list[str] = []
            with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
                reader = csv.reader(handle, delimiter=delimiter)
                for row in reader:
                    text = " | ".join(str(item).strip() for item in row if str(item).strip())
                    if text:
                        lines.append(text)
            return "\n".join(lines).strip()

        if suffix in TEXT_INDEX_EXTENSIONS:
            text = _read_text_fallback(path)
            if suffix == ".html":
                return _clean_html(text)
            return text.strip()

        if suffix == ".docx":
            return _extract_docx_text(path)
        if suffix == ".pdf":
            return _extract_pdf_text(path)
        if suffix == ".xlsx":
            return _extract_xlsx_text(path)
        if suffix == ".pptx":
            return _extract_pptx_text(path)
    except Exception:
        return ""
    return ""


def get_files_tree(settings: Settings, *, include_disabled_roots: bool = False) -> dict[str, Any]:
    roots = list_files_roots(settings, include_disabled=include_disabled_roots)
    nodes: list[dict[str, Any]] = []
    all_ids: set[str] = set()
    children_map: dict[str, list[str]] = {}

    for root in roots:
        if not include_disabled_roots and not bool(root.get("enabled")):
            continue
        root_id = int(root["id"])
        root_path = Path(str(root["root_path"]))
        root_doc_id = files_folder_document_id(root_id, "")
        root_node = {
            "id": root_doc_id,
            "source": SOURCE_FILES,
            "title": str(root.get("name") or f"Root {root_id}"),
            "url": f"files://{root_id}/",
            "parentDocumentId": None,
            "type": "root",
            "children": [],
        }
        all_ids.add(root_doc_id)
        nodes.append(root_node)

        if not root_path.exists() or not root_path.is_dir():
            continue

        folder_nodes: dict[str, dict[str, Any]] = {"": root_node}
        for rel in _walk_root_files(root_path):
            rel_path = _normalize_rel_path(rel)
            if not rel_path:
                continue
            suffixes = _suffixes(rel_path)
            if suffixes and not suffixes.intersection(INDEXABLE_EXTENSIONS):
                continue

            parent = ""
            parts = rel_path.split("/")
            if len(parts) > 1:
                current = ""
                for part in parts[:-1]:
                    current = f"{current}/{part}" if current else part
                    if current in folder_nodes:
                        parent = current
                        continue
                    folder_id = files_folder_document_id(root_id, current)
                    parent_folder_id = files_folder_document_id(root_id, parent)
                    folder_node = {
                        "id": folder_id,
                        "source": SOURCE_FILES,
                        "title": part,
                        "url": f"files://{root_id}/{current}/",
                        "parentDocumentId": parent_folder_id,
                        "type": "folder",
                        "children": [],
                    }
                    folder_nodes[current] = folder_node
                    folder_nodes[parent]["children"].append(folder_node)
                    all_ids.add(folder_id)
                    children_map.setdefault(parent_folder_id, []).append(folder_id)
                    parent = current

            doc_id = files_document_id(root_id, rel_path)
            parent_id = files_folder_document_id(root_id, parent)
            file_node = {
                "id": doc_id,
                "source": SOURCE_FILES,
                "title": parts[-1],
                "url": f"files://{root_id}/{rel_path}",
                "parentDocumentId": parent_id,
                "type": "file",
                "children": [],
            }
            folder_nodes[parent]["children"].append(file_node)
            all_ids.add(doc_id)
            children_map.setdefault(parent_id, []).append(doc_id)

        def _sort(node: dict[str, Any]) -> None:
            kids = node.get("children")
            if not isinstance(kids, list) or not kids:
                return
            kids.sort(key=lambda item: (0 if item.get("type") in {"root", "folder"} else 1, str(item.get("title") or "").lower()))
            for child in kids:
                if isinstance(child, dict):
                    _sort(child)

        _sort(root_node)

    return {
        "nodes": nodes,
        "count": len(all_ids),
        "all_ids": all_ids,
        "children_map": {key: list(values) for key, values in children_map.items()},
    }


def collect_files_documents_for_sync(settings: Settings) -> tuple[list[dict[str, str]], list[str]]:
    roots = [root for root in list_files_roots(settings, include_disabled=False) if bool(root.get("enabled"))]
    if not roots:
        return [], []

    now = _utc_now()
    changed_docs: list[dict[str, str]] = []
    stale_document_ids: list[str] = []

    with _connect(settings.auth_db_path) as conn:
        existing_rows = conn.execute(
            "SELECT document_id, root_id, rel_path, mtime, size, hash, indexed_at FROM files_documents"
        ).fetchall()
        existing_by_id = {str(row["document_id"]): row for row in existing_rows}
        seen_ids: set[str] = set()

        for root in roots:
            root_id = int(root["id"])
            root_path = Path(str(root["root_path"]))
            if not root_path.exists() or not root_path.is_dir():
                continue

            for rel in _walk_root_files(root_path):
                rel_path = _normalize_rel_path(rel)
                if not rel_path:
                    continue
                suffixes = _suffixes(rel_path)
                if suffixes and not suffixes.intersection(INDEXABLE_EXTENSIONS):
                    continue
                if not suffixes and Path(rel_path).name.lower() not in {"dockerfile", "makefile"}:
                    continue

                full_path = _safe_path_join(root_path, rel_path)
                if not full_path.exists() or not full_path.is_file():
                    continue
                stat = full_path.stat()
                mtime = datetime.fromtimestamp(stat.st_mtime, tz=UTC).isoformat(timespec="seconds")
                size = int(stat.st_size)
                doc_id = files_document_id(root_id, rel_path)
                seen_ids.add(doc_id)
                prev = existing_by_id.get(doc_id)

                hash_value = ""
                changed = True
                if prev is not None:
                    prev_mtime = str(prev["mtime"] or "")
                    prev_size = int(prev["size"] or 0)
                    if prev_mtime == mtime and prev_size == size:
                        changed = False

                if changed:
                    hash_value = _file_hash_sha256(full_path)
                    if prev is not None and str(prev["hash"] or "") == hash_value:
                        changed = False
                elif prev is not None:
                    hash_value = str(prev["hash"] or "")

                text = ""
                if changed:
                    text = extract_file_text(full_path)
                    if text:
                        changed_docs.append(
                            {
                                "id": doc_id,
                                "source": SOURCE_FILES,
                                "title": Path(rel_path).name,
                                "text": text,
                                "url": f"files://{root_id}/{rel_path}",
                            }
                        )
                    else:
                        stale_document_ids.append(doc_id)

                conn.execute(
                    """
                    INSERT INTO files_documents (document_id, root_id, rel_path, mtime, size, hash, indexed_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                    ON CONFLICT(document_id) DO UPDATE SET
                      root_id=excluded.root_id,
                      rel_path=excluded.rel_path,
                      mtime=excluded.mtime,
                      size=excluded.size,
                      hash=excluded.hash,
                      indexed_at=excluded.indexed_at
                    """,
                    (
                        doc_id,
                        root_id,
                        rel_path,
                        mtime,
                        size,
                        hash_value or _file_hash_sha256(full_path),
                        now,
                    ),
                )

        stale = [doc_id for doc_id in existing_by_id if doc_id not in seen_ids]
        for document_id in stale:
            stale_document_ids.append(document_id)
        if stale_document_ids:
            conn.executemany(
                "DELETE FROM files_documents WHERE document_id=?",
                [(doc_id,) for doc_id in stale_document_ids],
            )

        conn.commit()

    return changed_docs, stale_document_ids


def get_files_document(settings: Settings, document_id: str) -> dict[str, Any] | None:
    parsed = parse_files_document_id(document_id)
    if parsed is None:
        return None
    root_id, rel_path, is_folder = parsed
    root = get_files_root(settings, root_id)
    if root is None:
        return None
    path = _safe_path_join(Path(str(root["root_path"])), rel_path)
    if is_folder:
        return {
            "document_id": files_folder_document_id(root_id, rel_path),
            "root_id": root_id,
            "root_name": str(root["name"]),
            "rel_path": rel_path,
            "is_folder": True,
            "exists": path.exists() and path.is_dir(),
        }
    if not path.exists() or not path.is_file():
        return None
    text = extract_file_text(path)
    return {
        "document_id": files_document_id(root_id, rel_path),
        "root_id": root_id,
        "root_name": str(root["name"]),
        "rel_path": rel_path,
        "is_folder": False,
        "title": path.name,
        "size": int(path.stat().st_size),
        "mime": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
        "text": text,
    }


def _upload_storage_root() -> Path:
    root = Path("/data/chat_uploads")
    _ensure_directory(root)
    return root


def _upsert_files_document_row(settings: Settings, *, root_id: int, rel_path: str, path: Path) -> None:
    if not path.exists() or not path.is_file():
        return
    clean_rel = _normalize_rel_path(rel_path)
    if not clean_rel:
        return
    doc_id = files_document_id(int(root_id), clean_rel)
    stat = path.stat()
    mtime = datetime.fromtimestamp(stat.st_mtime, tz=UTC).isoformat(timespec="seconds")
    size = int(stat.st_size)
    hash_value = _file_hash_sha256(path)
    now = _utc_now_precise()
    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            INSERT INTO files_documents (document_id, root_id, rel_path, mtime, size, hash, indexed_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(document_id) DO UPDATE SET
              root_id=excluded.root_id,
              rel_path=excluded.rel_path,
              mtime=excluded.mtime,
              size=excluded.size,
              hash=excluded.hash,
              indexed_at=excluded.indexed_at
            """,
            (doc_id, int(root_id), clean_rel, mtime, size, hash_value, now),
        )
        conn.commit()


def _delete_files_document_row(settings: Settings, document_id: str) -> None:
    clean_id = str(document_id or "").strip()
    if not clean_id:
        return
    with _connect(settings.auth_db_path) as conn:
        conn.execute("DELETE FROM files_documents WHERE document_id=?", (clean_id,))
        conn.commit()


def _artifact_storage_root() -> Path:
    root = Path("/data/chat_artifacts")
    _ensure_directory(root)
    return root


def _files_upload_policy(settings: Settings) -> dict[str, Any]:
    return get_files_feature_settings(settings)


def _can_use_indexable_upload(settings: Settings) -> bool:
    state = _files_upload_policy(settings)
    return bool(state.get("upload_allow_indexable", True))


def _upload_is_expired(expires_at: str | None) -> bool:
    expires = _parse_iso8601(expires_at)
    if expires is None:
        return False
    return expires <= datetime.now(UTC)


def _artifact_is_expired(expires_at: str | None) -> bool:
    expires = _parse_iso8601(expires_at)
    if expires is None:
        return False
    return expires <= datetime.now(UTC)


def create_chat_upload(
    settings: Settings,
    *,
    user_id: int,
    chat_id: str | None,
    mode: str,
    filename: str,
    content_type: str | None,
    data: bytes,
) -> dict[str, Any]:
    policy = _files_upload_policy(settings)
    clean_mode = str(mode or "ephemeral").strip().lower()
    if clean_mode not in {"ephemeral", "indexable"}:
        raise ValueError("mode must be 'ephemeral' or 'indexable'")
    if clean_mode == "indexable" and not _can_use_indexable_upload(settings):
        raise ValueError("indexable uploads are disabled by policy")

    clean_filename = _normalize_upload_filename(filename, fallback="upload")
    allowed_ext = {str(item).strip().lower() for item in policy.get("upload_allowed_extensions") or []}
    if not _is_ext_allowed(clean_filename, allowed_ext):
        raise ValueError("file extension is not allowed by upload policy")

    max_size_bytes = int(policy["upload_max_size_mb"]) * 1024 * 1024
    if len(data) <= 0:
        raise ValueError("uploaded file is empty")
    if len(data) > max_size_bytes:
        raise ValueError(f"uploaded file exceeds {policy['upload_max_size_mb']}MB")

    upload_id = uuid.uuid4().hex
    now_dt = datetime.now(UTC)
    expires_at = (now_dt + timedelta(hours=int(policy["upload_ttl_hours"]))).isoformat(timespec="seconds")

    user_root = _upload_storage_root() / f"user_{int(user_id)}"
    _ensure_directory(user_root)
    stored_name = f"{upload_id}_{_sanitize_name(clean_filename, fallback='upload')}"
    storage_path = user_root / stored_name
    storage_path.write_bytes(data)

    extracted_text = extract_file_text(storage_path)
    indexed_document_id: str | None = None

    if clean_mode == "indexable":
        roots = [row for row in list_files_roots(settings, include_disabled=False) if bool(row.get("enabled"))]
        if not roots:
            raise ValueError("indexable upload requires at least one enabled files root")
        root = roots[0]
        root_id = int(root["id"])
        root_path = Path(str(root["root_path"]))
        target_dir_rel = _normalize_rel_path(f"uploads/{int(user_id)}")
        target_dir = _safe_path_join(root_path, target_dir_rel)
        _ensure_directory(target_dir)
        target_filename = _next_available_filename(target_dir, clean_filename)
        target_rel = _normalize_rel_path(f"{target_dir_rel}/{target_filename}")
        target = _safe_path_join(root_path, target_rel)
        shutil.copy2(storage_path, target)
        indexed_document_id = files_document_id(root_id, target_rel)

        extracted_index_text = extract_file_text(target)
        if extracted_index_text.strip():
            try:
                upsert_document_to_rag(
                    settings,
                    {
                        "id": indexed_document_id,
                        "source": SOURCE_FILES,
                        "title": clean_filename,
                        "text": extracted_index_text,
                        "url": f"files://{root_id}/{target_rel}",
                    },
                )
            except Exception:
                pass
        if extracted_index_text and not extracted_text:
            extracted_text = extracted_index_text

    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            INSERT INTO chat_uploads (id, user_id, chat_id, mode, mime, size, filename, storage_path, indexed_document_id, extracted_text, expires_at, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                upload_id,
                int(user_id),
                str(chat_id or "").strip() or None,
                clean_mode,
                str(content_type or "").strip() or None,
                len(data),
                clean_filename,
                str(storage_path),
                indexed_document_id,
                extracted_text[:200_000] if extracted_text else "",
                expires_at,
                now_dt.isoformat(timespec="seconds"),
            ),
        )
        conn.commit()

    return {
        "id": upload_id,
        "user_id": int(user_id),
        "chat_id": str(chat_id or "").strip() or None,
        "mode": clean_mode,
        "mime": str(content_type or "").strip() or "application/octet-stream",
        "size": len(data),
        "filename": clean_filename,
        "expires_at": expires_at,
        "indexed_document_id": indexed_document_id,
        "text_preview": (extracted_text or "")[:2000],
    }


def _upload_row_to_dict(row: sqlite3.Row) -> dict[str, Any]:
    return {
        "id": str(row["id"]),
        "user_id": int(row["user_id"]),
        "chat_id": str(row["chat_id"] or "") or None,
        "mode": str(row["mode"]),
        "mime": str(row["mime"] or "") or "application/octet-stream",
        "size": int(row["size"]),
        "filename": str(row["filename"]),
        "storage_path": str(row["storage_path"]),
        "indexed_document_id": str(row["indexed_document_id"] or "") or None,
        "extracted_text": str(row["extracted_text"] or ""),
        "expires_at": str(row["expires_at"]),
        "created_at": str(row["created_at"]),
    }


def get_chat_upload(settings: Settings, upload_id: str) -> dict[str, Any] | None:
    with _connect(settings.auth_db_path) as conn:
        row = conn.execute(
            """
            SELECT id, user_id, chat_id, mode, mime, size, filename, storage_path, indexed_document_id, extracted_text, expires_at, created_at
            FROM chat_uploads
            WHERE id=?
            """,
            (str(upload_id or "").strip(),),
        ).fetchone()
    if row is None:
        return None
    data = _upload_row_to_dict(row)
    data["expired"] = _upload_is_expired(data.get("expires_at"))
    return data


def list_chat_uploads_by_ids(settings: Settings, upload_ids: list[str], *, owner_user_id: int) -> list[dict[str, Any]]:
    normalized = [str(item or "").strip() for item in upload_ids if str(item or "").strip()]
    if not normalized:
        return []
    placeholders = ",".join("?" for _ in normalized)
    with _connect(settings.auth_db_path) as conn:
        rows = conn.execute(
            f"""
            SELECT id, user_id, chat_id, mode, mime, size, filename, storage_path, indexed_document_id, extracted_text, expires_at, created_at
            FROM chat_uploads
            WHERE id IN ({placeholders}) AND user_id=?
            """,
            (*normalized, int(owner_user_id)),
        ).fetchall()
    result: list[dict[str, Any]] = []
    for row in rows:
        item = _upload_row_to_dict(row)
        if _upload_is_expired(item.get("expires_at")):
            continue
        result.append(item)
    return result


def delete_chat_upload(settings: Settings, upload_id: str) -> bool:
    current = get_chat_upload(settings, upload_id)
    if current is None:
        return False
    path = Path(str(current.get("storage_path") or ""))
    try:
        if path.exists() and path.is_file():
            path.unlink()
    except Exception:
        pass

    indexed_document_id = str(current.get("indexed_document_id") or "").strip()
    if indexed_document_id:
        parsed = parse_files_document_id(indexed_document_id)
        if parsed is not None:
            root_id, rel_path, is_folder = parsed
            if not is_folder and rel_path:
                root = get_files_root(settings, root_id)
                if root is not None:
                    root_path = Path(str(root.get("root_path") or ""))
                    try:
                        indexed_path = _safe_path_join(root_path, rel_path)
                        if indexed_path.exists() and indexed_path.is_file():
                            indexed_path.unlink()
                    except Exception:
                        pass

    with _connect(settings.auth_db_path) as conn:
        conn.execute("DELETE FROM chat_uploads WHERE id=?", (str(upload_id),))
        if indexed_document_id:
            conn.execute("DELETE FROM files_documents WHERE document_id=?", (indexed_document_id,))
        conn.commit()
    return True


def _render_artifact_docx(content: str, out_path: Path) -> None:
    try:
        from docx import Document
    except Exception as exc:
        raise ValueError("DOCX renderer is unavailable") from exc
    document = Document()
    for line in (content or "").splitlines() or [""]:
        document.add_paragraph(line)
    document.save(str(out_path))


def _render_artifact_xlsx(content: str, out_path: Path) -> None:
    try:
        from openpyxl import Workbook
    except Exception as exc:
        raise ValueError("XLSX renderer is unavailable") from exc

    workbook = Workbook()
    sheet: Any = workbook.active
    if sheet is None:
        sheet = workbook.create_sheet("Sheet1")
    sheet.title = "Sheet1"

    parsed: Any = None
    try:
        parsed = json.loads(content)
    except Exception:
        parsed = None

    if isinstance(parsed, list) and parsed and isinstance(parsed[0], dict):
        headers = sorted({str(key) for row in parsed for key in row.keys()})
        for col, header in enumerate(headers, start=1):
            sheet.cell(row=1, column=col, value=header)
        for row_idx, row in enumerate(parsed, start=2):
            for col, header in enumerate(headers, start=1):
                value = row.get(header)
                sheet.cell(row=row_idx, column=col, value="" if value is None else str(value))
    elif isinstance(parsed, list):
        for row_idx, row in enumerate(parsed, start=1):
            if isinstance(row, list):
                for col, value in enumerate(row, start=1):
                    sheet.cell(row=row_idx, column=col, value="" if value is None else str(value))
            else:
                sheet.cell(row=row_idx, column=1, value=str(row))
    else:
        for row_idx, line in enumerate((content or "").splitlines(), start=1):
            parts = [part.strip() for part in line.split(",")]
            for col, value in enumerate(parts, start=1):
                sheet.cell(row=row_idx, column=col, value=value)

    workbook.save(str(out_path))


def _render_artifact_pdf(content: str, out_path: Path) -> None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception as exc:
        raise ValueError("PDF renderer is unavailable") from exc

    page_width, page_height = A4
    c = canvas.Canvas(str(out_path), pagesize=A4)
    y = page_height - 40
    for raw_line in (content or "").splitlines() or [""]:
        line = raw_line
        while len(line) > 120:
            chunk = line[:120]
            c.drawString(40, y, chunk)
            y -= 16
            line = line[120:]
            if y < 40:
                c.showPage()
                y = page_height - 40
        c.drawString(40, y, line)
        y -= 16
        if y < 40:
            c.showPage()
            y = page_height - 40
    c.save()


def _render_artifact_text(content: str, out_path: Path) -> None:
    out_path.write_text(content or "", encoding="utf-8")


def _artifact_extension(fmt: str) -> str:
    mapping = {
        "docx": ".docx",
        "xlsx": ".xlsx",
        "pdf": ".pdf",
        "md": ".md",
        "txt": ".txt",
        "csv": ".csv",
        "json": ".json",
    }
    return mapping.get(fmt, ".txt")


def _render_artifact_file(format_name: str, content: str, out_path: Path) -> None:
    fmt = str(format_name or "").strip().lower()
    if fmt == "docx":
        _render_artifact_docx(content, out_path)
        return
    if fmt == "xlsx":
        _render_artifact_xlsx(content, out_path)
        return
    if fmt == "pdf":
        _render_artifact_pdf(content, out_path)
        return
    if fmt in {"md", "txt", "csv", "json"}:
        _render_artifact_text(content, out_path)
        return
    raise ValueError("Unsupported artifact format")


def _create_artifact_token(settings: Settings, artifact_id: str, *, ttl_hours: int = 24) -> str:
    token = secrets.token_urlsafe(32)
    now = datetime.now(UTC)
    expires_at = (now + timedelta(hours=max(1, ttl_hours))).isoformat(timespec="seconds")
    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            "INSERT INTO chat_artifact_tokens (token, artifact_id, expires_at, used_at, created_at) VALUES (?, ?, ?, NULL, ?)",
            (token, str(artifact_id), expires_at, now.isoformat(timespec="seconds")),
        )
        conn.commit()
    return token


def create_artifact(
    settings: Settings,
    *,
    user_id: int,
    chat_id: str | None,
    format_name: str,
    filename: str | None,
    content: str,
) -> dict[str, Any]:
    policy = get_files_feature_settings(settings)
    fmt = str(format_name or "").strip().lower()
    if fmt not in {"docx", "xlsx", "pdf", "md", "txt", "csv", "json"}:
        raise ValueError("format must be one of: docx, xlsx, pdf, md, txt, csv, json")

    artifact_id = uuid.uuid4().hex
    extension = _artifact_extension(fmt)
    base = _sanitize_name(filename or f"artifact-{artifact_id}", fallback=f"artifact-{artifact_id}")
    if not base.lower().endswith(extension):
        base = f"{Path(base).stem}{extension}"

    user_root = _artifact_storage_root() / f"user_{int(user_id)}"
    _ensure_directory(user_root)
    out_path = user_root / f"{artifact_id}_{base}"
    _render_artifact_file(fmt, content, out_path)

    size = int(out_path.stat().st_size)
    max_bytes = int(policy["artifact_max_size_mb"]) * 1024 * 1024
    if size > max_bytes:
        out_path.unlink(missing_ok=True)
        raise ValueError(f"artifact exceeds {policy['artifact_max_size_mb']}MB")

    sha256 = _file_hash_sha256(out_path)
    now_dt = datetime.now(UTC)
    expires_at = (now_dt + timedelta(days=int(policy["artifact_ttl_days"]))).isoformat(timespec="seconds")

    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            INSERT INTO chat_artifacts (id, user_id, chat_id, format, filename, storage_path, size, sha256, expires_at, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                artifact_id,
                int(user_id),
                str(chat_id or "").strip() or None,
                fmt,
                base,
                str(out_path),
                size,
                sha256,
                expires_at,
                now_dt.isoformat(timespec="seconds"),
            ),
        )
        conn.commit()

    token = _create_artifact_token(settings, artifact_id)
    return {
        "id": artifact_id,
        "user_id": int(user_id),
        "chat_id": str(chat_id or "").strip() or None,
        "format": fmt,
        "filename": base,
        "size": size,
        "sha256": sha256,
        "expires_at": expires_at,
        "download_url": f"/assistant/artifacts/{artifact_id}/download?token={token}",
    }


def _artifact_row_to_dict(row: sqlite3.Row) -> dict[str, Any]:
    return {
        "id": str(row["id"]),
        "user_id": int(row["user_id"]),
        "chat_id": str(row["chat_id"] or "") or None,
        "format": str(row["format"]),
        "filename": str(row["filename"]),
        "storage_path": str(row["storage_path"]),
        "size": int(row["size"]),
        "sha256": str(row["sha256"]),
        "expires_at": str(row["expires_at"]),
        "created_at": str(row["created_at"]),
    }


def get_artifact(settings: Settings, artifact_id: str) -> dict[str, Any] | None:
    with _connect(settings.auth_db_path) as conn:
        row = conn.execute(
            """
            SELECT id, user_id, chat_id, format, filename, storage_path, size, sha256, expires_at, created_at
            FROM chat_artifacts
            WHERE id=?
            """,
            (str(artifact_id or "").strip(),),
        ).fetchone()
    if row is None:
        return None
    data = _artifact_row_to_dict(row)
    data["expired"] = _artifact_is_expired(data.get("expires_at"))
    return data


def issue_artifact_download_token(settings: Settings, artifact_id: str, *, ttl_hours: int = 24) -> str:
    return _create_artifact_token(settings, artifact_id, ttl_hours=ttl_hours)


def verify_artifact_token(settings: Settings, artifact_id: str, token: str) -> bool:
    clean_token = str(token or "").strip()
    if not clean_token:
        return False
    now = datetime.now(UTC)
    with _connect(settings.auth_db_path) as conn:
        row = conn.execute(
            "SELECT token, artifact_id, expires_at, used_at FROM chat_artifact_tokens WHERE token=?",
            (clean_token,),
        ).fetchone()
        if row is None:
            return False
        if str(row["artifact_id"] or "") != str(artifact_id):
            return False
        if str(row["used_at"] or "").strip():
            return False
        expires = _parse_iso8601(str(row["expires_at"] or ""))
        if expires is None or expires <= now:
            return False
        conn.execute(
            "UPDATE chat_artifact_tokens SET used_at=? WHERE token=?",
            (_utc_now_precise(), clean_token),
        )
        conn.commit()
    return True


def delete_artifact(settings: Settings, artifact_id: str) -> bool:
    artifact = get_artifact(settings, artifact_id)
    if artifact is None:
        return False
    path = Path(str(artifact.get("storage_path") or ""))
    try:
        if path.exists() and path.is_file():
            path.unlink()
    except Exception:
        pass
    with _connect(settings.auth_db_path) as conn:
        conn.execute("DELETE FROM chat_artifacts WHERE id=?", (str(artifact_id),))
        conn.execute("DELETE FROM chat_artifact_tokens WHERE artifact_id=?", (str(artifact_id),))
        conn.commit()
    return True


def cleanup_expired_assets(settings: Settings) -> dict[str, int]:
    now = datetime.now(UTC)
    deleted_uploads = 0
    deleted_artifacts = 0

    with _connect(settings.auth_db_path) as conn:
        upload_rows = conn.execute(
            "SELECT id, storage_path FROM chat_uploads WHERE expires_at <= ?",
            (now.isoformat(timespec="seconds"),),
        ).fetchall()
        for row in upload_rows:
            path = Path(str(row["storage_path"] or ""))
            try:
                if path.exists() and path.is_file():
                    path.unlink()
            except Exception:
                pass
            deleted_uploads += 1
        if upload_rows:
            conn.execute("DELETE FROM chat_uploads WHERE expires_at <= ?", (now.isoformat(timespec="seconds"),))

        artifact_rows = conn.execute(
            "SELECT id, storage_path FROM chat_artifacts WHERE expires_at <= ?",
            (now.isoformat(timespec="seconds"),),
        ).fetchall()
        for row in artifact_rows:
            path = Path(str(row["storage_path"] or ""))
            try:
                if path.exists() and path.is_file():
                    path.unlink()
            except Exception:
                pass
            deleted_artifacts += 1
        if artifact_rows:
            conn.execute("DELETE FROM chat_artifacts WHERE expires_at <= ?", (now.isoformat(timespec="seconds"),))

        conn.execute(
            "DELETE FROM chat_artifact_tokens WHERE used_at IS NOT NULL OR expires_at <= ?",
            (now.isoformat(timespec="seconds"),),
        )
        conn.commit()

    return {"deleted_uploads": deleted_uploads, "deleted_artifacts": deleted_artifacts}


def list_files_audit(settings: Settings, *, limit: int = 200, user_id: int | None = None) -> list[dict[str, Any]]:
    clean_limit = max(1, min(int(limit), 1000))
    with _connect(settings.auth_db_path) as conn:
        if user_id is None:
            rows = conn.execute(
                """
                SELECT id, user_id, op, path, document_id, before_hash, after_hash, status, message, created_at
                FROM files_write_audit
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (clean_limit,),
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT id, user_id, op, path, document_id, before_hash, after_hash, status, message, created_at
                FROM files_write_audit
                WHERE user_id=?
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (int(user_id), clean_limit),
            ).fetchall()
    return [
        {
            "id": str(row["id"]),
            "user_id": int(row["user_id"]),
            "op": str(row["op"]),
            "path": str(row["path"] or "") or None,
            "document_id": str(row["document_id"] or "") or None,
            "before_hash": str(row["before_hash"] or "") or None,
            "after_hash": str(row["after_hash"] or "") or None,
            "status": str(row["status"]),
            "message": str(row["message"] or "") or None,
            "created_at": str(row["created_at"]),
        }
        for row in rows
    ]


def _diff_text(before: str, after: str, path_label: str) -> str:
    before_lines = (before or "").splitlines()
    after_lines = (after or "").splitlines()
    diff = difflib.unified_diff(
        before_lines,
        after_lines,
        fromfile=f"a/{path_label}",
        tofile=f"b/{path_label}",
        lineterm="",
    )
    return "\n".join(diff)


def _path_is_text_writable(path: str) -> bool:
    suffixes = _suffixes(path)
    return bool(suffixes.intersection(TEXT_WRITE_EXTENSIONS))


def _write_audit_insert(
    settings: Settings,
    *,
    user_id: int,
    op: str,
    path: str | None,
    document_id: str | None,
    before_hash: str | None,
    after_hash: str | None,
    status: str,
    message: str | None,
) -> str:
    audit_id = uuid.uuid4().hex
    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            INSERT INTO files_write_audit (id, user_id, op, path, document_id, before_hash, after_hash, status, message, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                audit_id,
                int(user_id),
                str(op),
                str(path or "") or None,
                str(document_id or "") or None,
                str(before_hash or "") or None,
                str(after_hash or "") or None,
                str(status),
                str(message or "") or None,
                _utc_now_precise(),
            ),
        )
        conn.commit()
    return audit_id


def _write_audit_get(settings: Settings, audit_id: str) -> dict[str, Any] | None:
    with _connect(settings.auth_db_path) as conn:
        row = conn.execute(
            """
            SELECT id, user_id, op, path, document_id, before_hash, after_hash, status, message, created_at
            FROM files_write_audit
            WHERE id=?
            """,
            (str(audit_id or "").strip(),),
        ).fetchone()
    if row is None:
        return None
    return {
        "id": str(row["id"]),
        "user_id": int(row["user_id"]),
        "op": str(row["op"]),
        "path": str(row["path"] or "") or None,
        "document_id": str(row["document_id"] or "") or None,
        "before_hash": str(row["before_hash"] or "") or None,
        "after_hash": str(row["after_hash"] or "") or None,
        "status": str(row["status"]),
        "message": str(row["message"] or "") or None,
        "created_at": str(row["created_at"]),
    }


def _write_audit_update_status(
    settings: Settings,
    audit_id: str,
    *,
    status: str,
    before_hash: str | None = None,
    after_hash: str | None = None,
    message: str | None = None,
) -> None:
    with _connect(settings.auth_db_path) as conn:
        conn.execute(
            """
            UPDATE files_write_audit
            SET status=?, before_hash=COALESCE(?, before_hash), after_hash=COALESCE(?, after_hash), message=COALESCE(?, message)
            WHERE id=?
            """,
            (
                str(status),
                str(before_hash or "") or None,
                str(after_hash or "") or None,
                str(message or "") or None,
                str(audit_id),
            ),
        )
        conn.commit()


def _write_payload(audit_payload: dict[str, Any]) -> str:
    return json.dumps(audit_payload, ensure_ascii=False)


def _parse_write_payload(raw: str | None) -> dict[str, Any]:
    text = str(raw or "").strip()
    if not text:
        return {}
    try:
        data = json.loads(text)
    except json.JSONDecodeError:
        return {}
    return data if isinstance(data, dict) else {}


def preview_write_operation(
    settings: Settings,
    *,
    user_id: int,
    op: str,
    root_id: int,
    path: str,
    content: str | None = None,
    new_path: str | None = None,
) -> dict[str, Any]:
    if files_access_mode(settings) != FILES_ACCESS_READ_WRITE:
        raise ValueError("files source is in read-only mode")

    clean_op = str(op or "").strip().lower()
    if clean_op not in {"update", "create", "move", "delete"}:
        raise ValueError("op must be one of: update, create, move, delete")

    root = get_files_root(settings, int(root_id))
    if root is None or not bool(root.get("enabled")):
        raise ValueError("root not found or disabled")
    root_path = Path(str(root["root_path"]))
    if not root_path.exists() or not root_path.is_dir():
        raise ValueError("root path does not exist")

    rel_path = _normalize_rel_path(path)
    if not rel_path:
        raise ValueError("path is required")
    if clean_op in {"create", "update"} and not _path_is_text_writable(rel_path):
        raise ValueError("path extension is not writable in v1.1")

    target = _safe_path_join(root_path, rel_path)
    before = target.read_text(encoding="utf-8", errors="ignore") if target.exists() and target.is_file() else ""
    before_hash = hashlib.sha256(before.encode("utf-8")).hexdigest() if before else None

    payload: dict[str, Any] = {
        "op": clean_op,
        "root_id": int(root_id),
        "path": rel_path,
    }

    if clean_op == "create" and target.exists():
        raise ValueError("file already exists")
    if clean_op in {"update", "delete", "move"} and not target.exists():
        raise ValueError("target path does not exist")
    if clean_op in {"create", "update"}:
        after = str(content or "")
        payload["content"] = after
    elif clean_op == "move":
        rel_new = _normalize_rel_path(str(new_path or ""))
        if not rel_new:
            raise ValueError("new_path is required for move")
        destination = _safe_path_join(root_path, rel_new)
        if destination.exists():
            raise ValueError("new_path already exists")
        payload["new_path"] = rel_new
        after = before
    else:
        after = ""

    if clean_op == "move":
        diff = _diff_text(before, before, rel_path)
    elif clean_op == "delete":
        diff = _diff_text(before, "", rel_path)
    else:
        diff = _diff_text(before, after, rel_path)
    after_hash = hashlib.sha256(after.encode("utf-8")).hexdigest() if after else None

    doc_id = files_document_id(int(root_id), rel_path)
    audit_id = _write_audit_insert(
        settings,
        user_id=int(user_id),
        op=clean_op,
        path=rel_path,
        document_id=doc_id,
        before_hash=before_hash,
        after_hash=after_hash,
        status="preview",
        message=_write_payload(payload),
    )

    return {
        "audit_id": audit_id,
        "op": clean_op,
        "root_id": int(root_id),
        "path": rel_path,
        "document_id": doc_id,
        "diff": diff,
        "requires_confirm": True,
    }


def apply_write_operation(settings: Settings, *, audit_id: str, user_id: int) -> dict[str, Any]:
    if files_access_mode(settings) != FILES_ACCESS_READ_WRITE:
        raise ValueError("files source is in read-only mode")

    audit = _write_audit_get(settings, audit_id)
    if audit is None:
        raise ValueError("preview not found")
    if str(audit.get("status") or "") != "preview":
        raise ValueError("preview is already applied or invalid")
    if int(audit.get("user_id") or 0) != int(user_id):
        raise ValueError("cannot apply another user's preview")

    payload = _parse_write_payload(audit.get("message"))
    op = str(payload.get("op") or "").strip().lower()
    root_id = int(payload.get("root_id") or 0)
    rel_path = _normalize_rel_path(str(payload.get("path") or ""))
    if not root_id or not rel_path:
        raise ValueError("invalid preview payload")
    root = get_files_root(settings, root_id)
    if root is None or not bool(root.get("enabled")):
        raise ValueError("root not found or disabled")
    root_path = Path(str(root["root_path"]))
    target = _safe_path_join(root_path, rel_path)

    try:
        if op in {"create", "update"}:
            content = str(payload.get("content") or "")
            _ensure_directory(target.parent)
            target.write_text(content, encoding="utf-8")
            after_hash = hashlib.sha256(content.encode("utf-8")).hexdigest() if content else None
            _upsert_files_document_row(settings, root_id=root_id, rel_path=rel_path, path=target)
            doc_id = files_document_id(root_id, rel_path)
            indexed_text = extract_file_text(target)
            if indexed_text.strip():
                try:
                    upsert_document_to_rag(
                        settings,
                        {
                            "id": doc_id,
                            "source": SOURCE_FILES,
                            "title": target.name,
                            "text": indexed_text,
                            "url": f"files://{root_id}/{rel_path}",
                        },
                    )
                except Exception:
                    pass
            else:
                try:
                    delete_document_from_rag(settings, doc_id)
                except Exception:
                    pass
            _write_audit_update_status(settings, audit_id, status="applied", after_hash=after_hash)
            return {
                "ok": True,
                "audit_id": audit_id,
                "op": op,
                "document_id": doc_id,
            }

        if op == "move":
            rel_new = _normalize_rel_path(str(payload.get("new_path") or ""))
            if not rel_new:
                raise ValueError("new_path is missing")
            destination = _safe_path_join(root_path, rel_new)
            _ensure_directory(destination.parent)
            shutil.move(str(target), str(destination))
            old_doc_id = files_document_id(root_id, rel_path)
            new_doc_id = files_document_id(root_id, rel_new)
            _delete_files_document_row(settings, old_doc_id)
            _upsert_files_document_row(settings, root_id=root_id, rel_path=rel_new, path=destination)
            try:
                delete_document_from_rag(settings, old_doc_id)
            except Exception:
                pass
            indexed_text = extract_file_text(destination)
            if indexed_text.strip():
                try:
                    upsert_document_to_rag(
                        settings,
                        {
                            "id": new_doc_id,
                            "source": SOURCE_FILES,
                            "title": destination.name,
                            "text": indexed_text,
                            "url": f"files://{root_id}/{rel_new}",
                        },
                    )
                except Exception:
                    pass
            _write_audit_update_status(settings, audit_id, status="applied")
            return {
                "ok": True,
                "audit_id": audit_id,
                "op": op,
                "document_id": new_doc_id,
            }

        if op == "delete":
            trash_root = root_path / ".trash"
            _ensure_directory(trash_root)
            trash_name = f"{datetime.now(UTC).strftime('%Y%m%d%H%M%S')}_{target.name}"
            trash_path = trash_root / trash_name
            shutil.move(str(target), str(trash_path))
            doc_id = files_document_id(root_id, rel_path)
            _delete_files_document_row(settings, doc_id)
            try:
                delete_document_from_rag(settings, doc_id)
            except Exception:
                pass
            _write_audit_update_status(settings, audit_id, status="applied")
            return {
                "ok": True,
                "audit_id": audit_id,
                "op": op,
                "document_id": doc_id,
            }
    except Exception as exc:
        _write_audit_update_status(settings, audit_id, status="failed", message=str(exc))
        raise ValueError(str(exc)) from exc

    raise ValueError("unsupported preview operation")
